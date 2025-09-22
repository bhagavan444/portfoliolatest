from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import os
import uuid
from datetime import datetime
import requests
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from threading import Lock

app = Flask(__name__)
CORS(app)

# Thread-safe in-memory storage
chat_sessions = {}
chat_lock = Lock()

# Gemini API configuration
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "AIzaSyA3uC2IDU_Rb8VoJ-k2lGILQROc7j0SgNU")
GEMINI_API_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent"

# Downloads folder
os.makedirs('downloads', exist_ok=True)

# -------------------- Chat Endpoints --------------------
@app.route("/api/chats", methods=["GET"])
def get_chats():
    with chat_lock:
        sessions = [{"_id": k, "messages": v["messages"], "title": v.get("title", "Untitled Chat")}
                    for k, v in chat_sessions.items()]
    return jsonify({"sessions": sessions})

@app.route("/api/chats/<chat_id>", methods=["GET"])
def get_chat(chat_id):
    with chat_lock:
        session = chat_sessions.get(chat_id)
    if session:
        return jsonify(session)
    return jsonify({"error": "Chat session not found"}), 404

@app.route("/api/chats", methods=["POST"])
def create_chat():
    data = request.get_json()
    user_input = data.get("message", "")
    reply = data.get("reply", "")
    chat_id = data.get("chat_id", str(uuid.uuid4()))
    title = data.get("title")
    if not title:
        first_line = user_input.strip().split("\n")[0]
        title = first_line[:30] + "..." if len(first_line) > 30 else first_line or f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    with chat_lock:
        if chat_id not in chat_sessions:
            chat_sessions[chat_id] = {"messages": [], "title": title}
    return jsonify({"chat_id": chat_id})

@app.route("/api/chats/<chat_id>", methods=["PATCH"])
def update_chat(chat_id):
    data = request.get_json()
    new_title = data.get("title")
    if not new_title:
        return jsonify({"error": "Title is required"}), 400
    with chat_lock:
        if chat_id in chat_sessions:
            chat_sessions[chat_id]["title"] = new_title
            return jsonify({"message": "Title updated successfully", "title": new_title})
    return jsonify({"error": "Chat session not found"}), 404

@app.route("/api/chats/<chat_id>", methods=["DELETE"])
def delete_chat(chat_id):
    with chat_lock:
        if chat_id in chat_sessions:
            del chat_sessions[chat_id]
            return jsonify({"message": "Chat deleted successfully"})
    return jsonify({"error": "Chat session not found"}), 404

@app.route("/api/chats", methods=["DELETE"])
def delete_all_chats():
    with chat_lock:
        chat_sessions.clear()
    return jsonify({"message": "All chats deleted successfully"})

# -------------------- Per-message deletion --------------------
@app.route("/api/chats/<chat_id>/message/<msg_id>", methods=["DELETE"])
def delete_message(chat_id, msg_id):
    with chat_lock:
        if chat_id in chat_sessions:
            messages = chat_sessions[chat_id]["messages"]
            chat_sessions[chat_id]["messages"] = [m for m in messages if m["id"] != msg_id]
            return jsonify({"message": "Message deleted successfully"})
    return jsonify({"error": "Chat or message not found"}), 404

# -------------------- Chat search/filter --------------------
@app.route("/api/chats/search", methods=["GET"])
def search_chats():
    query = request.args.get("q", "").lower()
    with chat_lock:
        results = [
            {"_id": k, "title": v["title"]}
            for k, v in chat_sessions.items()
            if query in v["title"].lower() or any(query in (m.get("message") or "").lower() for m in v["messages"])
        ]
    return jsonify({"results": results})

# -------------------- File Download --------------------
@app.route("/download/<filename>", methods=["GET"])
def download(filename):
    safe_path = os.path.join("downloads", filename)
    if os.path.exists(safe_path):
        return send_from_directory('downloads', filename, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

# -------------------- Chat + File Upload & ATS --------------------
@app.route("/api/chat", methods=["POST"])
def chat():
    try:
        user_input, file, chat_id = "", None, None
        if request.content_type.startswith("multipart/form-data"):
            user_input = request.form.get("message", "")
            chat_id = request.form.get("chat_id")
            file = request.files.get("file")
        else:
            data = request.get_json()
            user_input = data.get("message", "")
            chat_id = data.get("chat_id")

        if not user_input and not file:
            return jsonify({"reply": "⚠️ No input or file received."}), 400

        # Process file
        filename, ext, file_text = None, None, ""
        if file:
            filename = f"{uuid.uuid4().hex}_{file.filename}"
            filepath = os.path.join("downloads", filename)
            file.save(filepath)
            ext = os.path.splitext(filename)[1].lower()
            if ext == ".pdf":
                reader = PdfReader(filepath)
                file_text = "\n".join(page.extract_text() or "" for page in reader.pages)
            elif ext == ".docx":
                doc = Document(filepath)
                file_text = "\n".join([p.text for p in doc.paragraphs])
            elif ext == ".pptx":
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            file_text += shape.text + "\n"
            elif ext in [".jpg", ".jpeg", ".png"]:
                file_text = "[Image uploaded, text analysis not available]"
            else:
                return jsonify({"reply": "⚠️ Unsupported file type."}), 400
            user_input += f"\n\nFile Content ({filename}): {file_text}"

        # Auto-detect resume
        resume_keywords = ["experience", "education", "skills", "projects", "certifications", "objective", "summary"]
        is_resume = any(word.lower() in user_input.lower() for word in resume_keywords)

        # Gemini prompt
        system_prompt = (
            "You are an expert ATS analyzer.\nAnalyze this resume content and provide an ATS score (0-100) with detailed improvement suggestions in plain text."
            if is_resume else
            "You are a helpful assistant.\nSummarize the content or answer user's query in plain text."
        )
        prompt_text = system_prompt + "\n\n" + user_input

        headers = {"Content-Type": "application/json", "x-goog-api-key": GEMINI_API_KEY}
        data = {"contents": [{"role": "user", "parts": [{"text": prompt_text}]}]}
        response = requests.post(GEMINI_API_URL, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        reply = result.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")

        # Store chat messages (continuous)
        new_chat_id = chat_id or str(uuid.uuid4())
        timestamp = datetime.now().strftime("%I:%M %p")
        with chat_lock:
            if new_chat_id not in chat_sessions:
                chat_sessions[new_chat_id] = {"messages": [], "title": user_input[:20] + "..."}

            # Append user message
            chat_sessions[new_chat_id]["messages"].append({
                "id": str(uuid.uuid4()),
                "message": user_input,
                "reply": None,
                "time": timestamp,
                "role": "user",
                "file": {"name": filename} if file else None
            })

            # Append assistant reply
            chat_sessions[new_chat_id]["messages"].append({
                "id": str(uuid.uuid4()),
                "message": None,
                "reply": reply,
                "time": timestamp,
                "role": "assistant"
            })

        return jsonify({"reply": reply, "chat_id": new_chat_id})

    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({"reply": "⚠️ Error processing your message."}), 500

if __name__ == "__main__":
    app.run(debug=True, port=5000, threaded=True)
